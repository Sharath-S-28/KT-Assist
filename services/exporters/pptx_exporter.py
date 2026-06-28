"""
services/exporters/pptx_exporter.py — KT Assurance Report PowerPoint export
(Phase 10 / Session 32).

Renders schemas.assurance_report.AssuranceReport to a .pptx, one slide per
Appendix C section, using python-pptx. Library reconciliation: the pptx
skill's pptxgenjs guidance targets interactively-authored, visually
designed decks built turn-by-turn in conversation; this is a backend
FastAPI export endpoint producing a data-driven artifact from a Pydantic
model with no design brief, so python-pptx (a Python library this
service layer can call directly, with no Node toolchain in the request
path) is the appropriate tool, mirroring the same reconciliation the PDF
exporter makes for reportlab over the pdf skill's other options. Like
every other Phase 10 module, this is a formatter -- it touches no score
and adds no new numbers.
"""

from pptx import Presentation
from pptx.util import Inches, Pt

from schemas.assurance_report import AssuranceReport

_TITLE_LAYOUT = 0
_CONTENT_LAYOUT = 1


def _add_bullet_slide(prs: Presentation, title: str, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[_CONTENT_LAYOUT])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    if not bullets:
        body.text = "No data available."
        return
    body.text = bullets[0]
    for line in bullets[1:]:
        paragraph = body.add_paragraph()
        paragraph.text = line
        paragraph.font.size = Pt(16)


def export_assurance_report_pptx(report: AssuranceReport, output_path: str) -> str:
    prs = Presentation()

    # 1. Cover
    cover = prs.slides.add_slide(prs.slide_layouts[_TITLE_LAYOUT])
    cover.shapes.title.text = "KT Assurance Report"
    cover.placeholders[1].text = (
        f"{report.program_name}\nReport ID: {report.report_id}\n"
        f"Generated: {report.generated_at.isoformat()}"
    )

    # 2. Executive summary
    health = report.program_health
    _add_bullet_slide(
        prs,
        "Executive Summary",
        [
            f"Lifecycle State: {health.lifecycle_state}",
            f"Completion Status: {health.completion_status}",
            f"Readiness: {health.readiness or 'Not Assessed'}",
            f"At Risk: {'Yes' if health.at_risk else 'No'}",
            f"Average Coverage: {f'{health.coverage:.0%}' if health.coverage is not None else 'N/A'}",
            f"Average OIS: {f'{health.ois:.1f}' if health.ois is not None else 'N/A'}",
            f"Overall Decision: {report.overall_decision or 'Not Assessed'}",
        ],
    )

    # 3. Coverage assessment
    _add_bullet_slide(
        prs,
        "Coverage Assessment",
        [
            f"{cov.package_id}: {cov.coverage:.0%} coverage, "
            f"{'sufficient' if cov.sufficient else 'insufficient'}"
            for cov in report.coverage_by_package
        ],
    )

    # 4. Gap & risk analysis
    gs = report.gap_summary
    gap_bullets = [
        f"Total gaps: {gs.total} (Open {gs.open}, Closed {gs.closed}, "
        f"Critical {gs.critical}, High risk {gs.high_risk})"
    ]
    gap_bullets += [
        f"{cell.domain}: {cell.open_gaps} open gaps, {cell.critical_gaps} critical"
        for cell in report.risk_concentration
    ]
    _add_bullet_slide(prs, "Gap & Risk Analysis", gap_bullets)

    # 5. Receiver readiness summary
    _add_bullet_slide(
        prs,
        "Receiver Readiness Summary",
        [
            f"{rd.receiver_name}: OIS {rd.ois:.1f}, {rd.readiness_status}, "
            f"certification {rd.certification or 'None'}"
            for rd in report.readiness_by_receiver
        ],
    )

    # 6. Competency assessment detail
    cs = report.competency_summary
    _add_bullet_slide(
        prs,
        "Competency Assessment Detail",
        [f"Total {cs.total} | Pass {cs.passed} | Fail {cs.failed} | Warning {cs.warning}"],
    )

    # 7. Certification & sign-off status
    _add_bullet_slide(
        prs,
        "Certification & Sign-off Status",
        [
            f"{cert.receiver_name}: {cert.readiness_status}, certification {cert.certification or 'None'}"
            for cert in report.certifications
        ],
    )

    # 8. Recommendations
    receiver_names = {rd.receiver_id: rd.receiver_name for rd in report.readiness_by_receiver}
    recommendation_bullets = []
    for receiver_id, items in report.recommendations_by_receiver.items():
        for item in items:
            recommendation_bullets.append(
                f"{receiver_names.get(receiver_id, receiver_id)} - {item.competency_name} "
                f"(score {item.score:.0f}): {'; '.join(item.actions)}"
            )
    _add_bullet_slide(prs, "Recommendations", recommendation_bullets)

    # 9. Traceability appendix
    _add_bullet_slide(
        prs,
        "Traceability Appendix",
        [f"Traced readiness record: {rid}" for rid in report.traced_receiver_readiness_ids],
    )

    prs.save(output_path)
    return output_path
