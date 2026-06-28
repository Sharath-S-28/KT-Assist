"""
services/asset_ingestion.py — Asset Ingestion & Text Extraction
(Phase 4 / KAI, Session 10).

Converts uploaded knowledge assets (PDF, DOCX, PPTX, TXT — meeting
notes, SOPs, runbooks, code/process docs all arrive as one of these
four container formats) into clean, chunked text. The raw file is
persisted to disk under config.ASSETS_DIR; only metadata (filename,
type, storage path, content hash, extraction status) is persisted to
SQLite via the KnowledgeAsset row — extracted text/chunks are derived
on demand from the stored file, never duplicated into the database.

KAI (Session 11) extracts knowledge objects from these chunks and
caches its output by content_hash, so this module's hashing is also
the cache key for downstream extraction cost control.
"""

import hashlib
import re
from pathlib import Path

import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
from sqlalchemy.orm import Session

import config
from models import KnowledgeAsset
from utils.errors import NotFoundError, ValidationFailedError

SUPPORTED_FILE_TYPES = {"pdf", "docx", "pptx", "txt"}

# Character-based chunking (scaffold values, recalibratable without
# changing the contract): large enough to preserve paragraph context
# for KAI extraction, small enough to stay well under typical prompt
# budgets. Overlap preserves cross-chunk context for relationship
# discovery (Session 12) at chunk boundaries.
CHUNK_SIZE_CHARS = 2000
CHUNK_OVERLAP_CHARS = 200


def _infer_file_type(filename: str) -> str:
    suffix = Path(filename).suffix.lower().lstrip(".")
    if suffix not in SUPPORTED_FILE_TYPES:
        raise ValidationFailedError(
            f"Unsupported file type {suffix!r}. Supported types: {sorted(SUPPORTED_FILE_TYPES)}.",
            details={"filename": filename, "suffix": suffix},
        )
    return suffix


def compute_content_hash(content: bytes) -> str:
    """SHA-256 hex digest of the raw file bytes — used both as a dedup
    key and as the KAI extraction cache key (Session 11)."""
    return hashlib.sha256(content).hexdigest()


def extract_raw_text(file_path: Path, file_type: str) -> str:
    """Dispatch text extraction by container format. Returns the raw,
    un-normalized text concatenated from every page/slide/paragraph."""
    if file_type == "txt":
        return file_path.read_text(encoding="utf-8", errors="replace")

    if file_type == "pdf":
        pages: list[str] = []
        with pdfplumber.open(str(file_path)) as pdf:
            for page in pdf.pages:
                pages.append(page.extract_text() or "")
        return "\n\n".join(pages)

    if file_type == "docx":
        doc = DocxDocument(str(file_path))
        return "\n".join(p.text for p in doc.paragraphs)

    if file_type == "pptx":
        prs = Presentation(str(file_path))
        slide_texts: list[str] = []
        for slide in prs.slides:
            shape_texts = [
                shape.text_frame.text
                for shape in slide.shapes
                if shape.has_text_frame
            ]
            slide_texts.append("\n".join(shape_texts))
        return "\n\n".join(slide_texts)

    raise ValidationFailedError(f"No extractor registered for file type {file_type!r}.")


def normalize_text(raw_text: str) -> str:
    """Collapse repeated whitespace/blank lines and strip trailing
    whitespace per line, without destroying paragraph boundaries
    (KAI's object extraction relies on paragraph structure)."""
    text = raw_text.replace("\r\n", "\n").replace("\r", "\n")
    lines = [line.strip() for line in text.split("\n")]
    text = "\n".join(lines)
    # Collapse 3+ consecutive blank lines down to a single blank line.
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def chunk_text(
    text: str,
    chunk_size: int = CHUNK_SIZE_CHARS,
    overlap: int = CHUNK_OVERLAP_CHARS,
) -> list[str]:
    """Split normalized text into overlapping character-window chunks.
    Prefers to break on a paragraph boundary near the window edge so
    chunks don't sever a sentence mid-word when avoidable."""
    text = text.strip()
    if not text:
        return []
    if len(text) <= chunk_size:
        return [text]
    if overlap >= chunk_size:
        raise ValidationFailedError("overlap must be smaller than chunk_size.")

    chunks: list[str] = []
    start = 0
    text_len = len(text)

    while start < text_len:
        end = min(start + chunk_size, text_len)
        if end < text_len:
            boundary = text.rfind("\n\n", start, end)
            if boundary != -1 and boundary > start:
                end = boundary
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= text_len:
            break
        start = max(end - overlap, start + 1)

    return chunks


def ingest_asset(
    db: Session,
    package_id: str,
    filename: str,
    content: bytes,
) -> tuple[KnowledgeAsset, list[str]]:
    """End-to-end ingestion: validate type, persist the raw file,
    create/refresh the KnowledgeAsset row, extract + normalize + chunk
    its text. Returns the asset row and its text chunks (not persisted
    — derived fresh from the stored file on every call)."""
    file_type = _infer_file_type(filename)
    content_hash = compute_content_hash(content)

    package_dir = config.ASSETS_DIR / package_id
    package_dir.mkdir(parents=True, exist_ok=True)

    asset = KnowledgeAsset(
        package_id=package_id,
        filename=filename,
        file_type=file_type,
        storage_path="",  # filled in below once we know the asset id
        content_hash=content_hash,
        extraction_status="Pending",
    )
    db.add(asset)
    db.flush()  # assigns asset.id

    storage_path = package_dir / f"{asset.id}_{filename}"
    storage_path.write_bytes(content)
    asset.storage_path = str(storage_path)

    try:
        raw_text = extract_raw_text(storage_path, file_type)
        normalized = normalize_text(raw_text)
        chunks = chunk_text(normalized)
        asset.extraction_status = "Extracted"
    except Exception:
        asset.extraction_status = "Failed"
        db.flush()
        raise

    db.flush()
    return asset, chunks


def get_asset_chunks(db: Session, asset_id: str) -> list[str]:
    """Re-derive chunks for an already-ingested asset by re-reading and
    re-extracting its stored file (chunks are never persisted)."""
    asset = db.query(KnowledgeAsset).filter_by(id=asset_id).first()
    if asset is None:
        raise NotFoundError(f"Knowledge asset {asset_id!r} not found.")

    raw_text = extract_raw_text(Path(asset.storage_path), asset.file_type)
    normalized = normalize_text(raw_text)
    return chunk_text(normalized)
