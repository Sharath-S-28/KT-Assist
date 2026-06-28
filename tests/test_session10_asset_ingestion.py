"""
tests/test_session10_asset_ingestion.py — Phase 4 / Session 10 success
criterion: mixed-format assets (PDF, DOCX, PPTX, TXT) are ingested and
converted into clean, chunked text.
"""

import io

import pytest
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.util import Inches
from reportlab.pdfgen import canvas

from services.asset_ingestion import (
    CHUNK_OVERLAP_CHARS,
    CHUNK_SIZE_CHARS,
    chunk_text,
    compute_content_hash,
    get_asset_chunks,
    ingest_asset,
    normalize_text,
)
from utils.errors import NotFoundError, ValidationFailedError


def _make_txt_bytes(text: str) -> bytes:
    return text.encode("utf-8")


def _make_docx_bytes(paragraphs: list[str]) -> bytes:
    doc = DocxDocument()
    for p in paragraphs:
        doc.add_paragraph(p)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_pptx_bytes(slide_texts: list[str]) -> bytes:
    prs = Presentation()
    layout = prs.slide_layouts[5]  # blank-ish layout with a title placeholder
    for text in slide_texts:
        slide = prs.slides.add_slide(layout)
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        box.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_pdf_bytes(lines: list[str]) -> bytes:
    buf = io.BytesIO()
    c = canvas.Canvas(buf)
    y = 750
    for line in lines:
        c.drawString(72, y, line)
        y -= 20
    c.save()
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Per-format ingestion
# ---------------------------------------------------------------------------

def test_ingest_txt_asset_extracts_clean_text(db_session, sample_package):
    content = _make_txt_bytes("Reconcile the GL.\n\nPost journal entries.")
    asset, chunks = ingest_asset(db_session, sample_package.id, "notes.txt", content)

    assert asset.file_type == "txt"
    assert asset.extraction_status == "Extracted"
    assert asset.content_hash == compute_content_hash(content)
    assert len(chunks) == 1
    assert "Reconcile the GL." in chunks[0]
    assert "Post journal entries." in chunks[0]


def test_ingest_docx_asset_extracts_paragraphs(db_session, sample_package):
    content = _make_docx_bytes(["Month-end close SOP.", "Step 1: reconcile GL.", "Step 2: post entries."])
    asset, chunks = ingest_asset(db_session, sample_package.id, "sop.docx", content)

    assert asset.file_type == "docx"
    assert asset.extraction_status == "Extracted"
    joined = "\n".join(chunks)
    assert "Month-end close SOP." in joined
    assert "Step 1: reconcile GL." in joined


def test_ingest_pptx_asset_extracts_slide_text(db_session, sample_package):
    content = _make_pptx_bytes(["Runbook overview", "Escalation contacts"])
    asset, chunks = ingest_asset(db_session, sample_package.id, "runbook.pptx", content)

    assert asset.file_type == "pptx"
    assert asset.extraction_status == "Extracted"
    joined = "\n".join(chunks)
    assert "Runbook overview" in joined
    assert "Escalation contacts" in joined


def test_ingest_pdf_asset_extracts_text(db_session, sample_package):
    content = _make_pdf_bytes(["Known issue log", "Issue 1: late batch job"])
    asset, chunks = ingest_asset(db_session, sample_package.id, "issues.pdf", content)

    assert asset.file_type == "pdf"
    assert asset.extraction_status == "Extracted"
    joined = "\n".join(chunks)
    assert "Known issue log" in joined


def test_ingest_rejects_unsupported_file_type(db_session, sample_package):
    with pytest.raises(ValidationFailedError):
        ingest_asset(db_session, sample_package.id, "diagram.png", b"not a real image")


def test_ingest_persists_metadata_and_storage_path(db_session, sample_package):
    content = _make_txt_bytes("Some content.")
    asset, _ = ingest_asset(db_session, sample_package.id, "doc.txt", content)

    from pathlib import Path
    assert Path(asset.storage_path).exists()
    assert Path(asset.storage_path).read_bytes() == content


def test_ingest_marks_failed_status_on_extraction_error(db_session, sample_package, monkeypatch):
    import services.asset_ingestion as mod

    def _boom(*args, **kwargs):
        raise RuntimeError("simulated extractor crash")

    monkeypatch.setattr(mod, "extract_raw_text", _boom)

    with pytest.raises(RuntimeError):
        mod.ingest_asset(db_session, sample_package.id, "doc.txt", b"hello")

    from models import KnowledgeAsset
    asset = db_session.query(KnowledgeAsset).filter_by(filename="doc.txt").first()
    assert asset.extraction_status == "Failed"


# ---------------------------------------------------------------------------
# Normalization
# ---------------------------------------------------------------------------

def test_normalize_text_collapses_blank_lines_and_trims_whitespace():
    raw = "Line one.   \r\n\r\n\r\n\r\nLine two.\r\n   \nLine three.  "
    normalized = normalize_text(raw)

    assert "\r" not in normalized
    assert "\n\n\n" not in normalized
    assert normalized.startswith("Line one.")
    assert normalized.endswith("Line three.")


# ---------------------------------------------------------------------------
# Chunking
# ---------------------------------------------------------------------------

def test_chunk_text_returns_single_chunk_for_short_text():
    text = "Short document."
    chunks = chunk_text(text)
    assert chunks == ["Short document."]


def test_chunk_text_returns_empty_list_for_empty_text():
    assert chunk_text("") == []
    assert chunk_text("   ") == []


def test_chunk_text_splits_long_text_into_overlapping_chunks():
    paragraph = "Sentence about the process. " * 50  # long single paragraph
    text = "\n\n".join([paragraph] * 5)  # well over CHUNK_SIZE_CHARS
    chunks = chunk_text(text, chunk_size=500, overlap=50)

    assert len(chunks) > 1
    for chunk in chunks:
        assert len(chunk) <= 500 + 50  # allow boundary slack from paragraph-snapping
    # Reassembled chunks should cover the original content (no big gaps).
    assert sum(len(c) for c in chunks) >= len(text) * 0.9


def test_chunk_text_rejects_overlap_larger_than_chunk_size():
    with pytest.raises(ValidationFailedError):
        chunk_text("x" * 100, chunk_size=50, overlap=60)


def test_chunk_text_default_constants_are_sane():
    assert CHUNK_OVERLAP_CHARS < CHUNK_SIZE_CHARS


# ---------------------------------------------------------------------------
# Re-derivation / caching key
# ---------------------------------------------------------------------------

def test_get_asset_chunks_rederives_same_chunks_as_ingestion(db_session, sample_package):
    content = _make_txt_bytes("Alpha paragraph.\n\nBeta paragraph.")
    asset, original_chunks = ingest_asset(db_session, sample_package.id, "doc.txt", content)

    rederived = get_asset_chunks(db_session, asset.id)
    assert rederived == original_chunks


def test_get_asset_chunks_raises_not_found_for_unknown_asset(db_session):
    with pytest.raises(NotFoundError):
        get_asset_chunks(db_session, "does-not-exist")


def test_identical_content_hashes_match_across_filenames(db_session, sample_package):
    content = _make_txt_bytes("Identical content.")
    asset1, _ = ingest_asset(db_session, sample_package.id, "a.txt", content)
    asset2, _ = ingest_asset(db_session, sample_package.id, "b.txt", content)

    assert asset1.content_hash == asset2.content_hash
